"""
Generador de presentación JISBD 2026
"UMDA aplicado a la Mejora Automática del Software"
Pablo Bermejo — UCLM, Albacete
"""

import io, os, textwrap
import cairosvg
import fitz  # pymupdf
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Colores ────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x38, 0x64)
BLUE2  = RGBColor(0x2E, 0x74, 0xB5)
ORANGE = RGBColor(0xC5, 0x5A, 0x11)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
DGRAY  = RGBColor(0x40, 0x40, 0x40)
GREEN  = RGBColor(0x37, 0x86, 0x44)
GOLD   = RGBColor(0xFF, 0xC0, 0x00)

# ── Rutas ──────────────────────────────────────────────────────────────────────
PLOTS = "/home/pablo/magpie/experiments/results_INTED2026/plots"
OUT   = "/home/pablo/magpie/experiments/results_INTED2026/slides/JISBD2026_Bermejo.pptx"

# ── Presentación base ──────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]   # completamente en blanco

W = prs.slide_width
H = prs.slide_height

# ── Helpers ────────────────────────────────────────────────────────────────────

def svg_to_png_bytes(path, scale=2):
    return cairosvg.svg2png(url=path, scale=scale)

def pdf_to_png_bytes(path, page=0, dpi=200):
    doc  = fitz.open(path)
    pg   = doc[page]
    mat  = fitz.Matrix(dpi / 72, dpi / 72)
    pix  = pg.get_pixmap(matrix=mat, alpha=False)
    return pix.tobytes("png")

def add_bg(slide, color=WHITE):
    bg   = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color:
        shp.line.color.rgb = line_color
    else:
        shp.line.fill.background()
    return shp

def txbox(slide, text, left, top, width, height,
          size=20, bold=False, color=DGRAY,
          align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_para(tf, text, size=20, bold=False, color=DGRAY,
             align=PP_ALIGN.LEFT, italic=False, space_before=0, indent=False):
    p   = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if indent:
        p.level = 1
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def header_bar(slide, title, sub=None, h=Cm(2.2)):
    bar = rect(slide, 0, 0, W, h, NAVY)
    tf  = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.5)
    tf.margin_top  = Cm(0.2)
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size  = Pt(28)
    run.font.bold  = True
    run.font.color.rgb = WHITE
    if sub:
        p2  = tf.add_paragraph()
        r2  = p2.add_run()
        r2.text = sub
        r2.font.size  = Pt(16)
        r2.font.color.rgb = GOLD
    return bar

def insert_image_bytes(slide, img_bytes, left, top, width=None, height=None):
    stream = io.BytesIO(img_bytes)
    if width and height:
        pic = slide.shapes.add_picture(stream, left, top, width, height)
    elif width:
        pic = slide.shapes.add_picture(stream, left, top, width=width)
    elif height:
        pic = slide.shapes.add_picture(stream, left, top, height=height)
    else:
        pic = slide.shapes.add_picture(stream, left, top)
    return pic

def bullet_frame(slide, items, left, top, width, height,
                 size=21, color=DGRAY, marker="▸ "):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        if isinstance(item, dict):
            run.text = marker + item["text"]
            run.font.size  = Pt(item.get("size", size))
            run.font.bold  = item.get("bold", False)
            run.font.color.rgb = item.get("color", color)
        else:
            run.text = marker + item
            run.font.size  = Pt(size)
            run.font.color.rgb = color
    return tb

def table_slide(slide, headers, rows,
                left, top, col_widths, row_height=Cm(0.9),
                head_fill=NAVY, head_color=WHITE,
                alt_fill=LGRAY, line_color=RGBColor(0xBF,0xBF,0xBF),
                font_size=18):
    """Draw a table manually using rectangles and text boxes."""
    n_cols = len(headers)
    n_rows = len(rows)
    # header row
    x = left
    for i, (hdr, cw) in enumerate(zip(headers, col_widths)):
        r = rect(slide, x, top, cw, row_height, head_fill, line_color)
        tb = slide.shapes.add_textbox(x + Cm(0.15), top + Cm(0.1),
                                      cw - Cm(0.3), row_height - Cm(0.1))
        tf = tb.text_frame
        p  = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.size  = Pt(font_size)
        run.font.bold  = True
        run.font.color.rgb = head_color
        x += cw
    # data rows
    for ri, row in enumerate(rows):
        y    = top + row_height * (ri + 1)
        fill = alt_fill if ri % 2 == 0 else WHITE
        x    = left
        for ci, (cell, cw) in enumerate(zip(row, col_widths)):
            r = rect(slide, x, y, cw, row_height, fill, line_color)
            tb = slide.shapes.add_textbox(x + Cm(0.15), y + Cm(0.05),
                                          cw - Cm(0.3), row_height)
            tf = tb.text_frame
            p  = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(cell)
            run.font.size = Pt(font_size)
            # highlight UMDA row
            if "UMDA" in str(row[0]):
                run.font.bold  = True
                run.font.color.rgb = NAVY
            elif "1" == str(cell) and ci > 0:
                run.font.bold = True
                run.font.color.rgb = GREEN
            else:
                run.font.color.rgb = DGRAY
            x += cw

def footer(slide, text="Pablo Bermejo · UCLM · JISBD 2026"):
    rect(slide, 0, H - Cm(0.7), W, Cm(0.7), NAVY)
    txbox(slide, text, Cm(0.5), H - Cm(0.65), W - Cm(2), Cm(0.6),
          size=11, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 1 — PORTADA
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl, NAVY)

# Franja decorativa central
rect(sl, 0, Cm(5.5), W, Cm(0.25), GOLD)

# Título
txbox(sl, "UMDA aplicado a la Mejora\nAutomática del Software",
      Cm(1.2), Cm(1.8), W - Cm(2.4), Cm(3.5),
      size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Autor + afiliación
txbox(sl, "Pablo Bermejo López",
      Cm(1.2), Cm(5.9), W - Cm(2.4), Cm(0.9),
      size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

txbox(sl, "Dpto. Sistemas Informáticos, SI y Minería de Datos (I3A)\nUniversidad de Castilla-La Mancha, Albacete",
      Cm(1.2), Cm(6.7), W - Cm(2.4), Cm(1.3),
      size=18, color=WHITE, align=PP_ALIGN.CENTER)

# Evento
rect(sl, Cm(3.5), Cm(8.5), Cm(26.4), Cm(0.85), BLUE2)
txbox(sl, "JISBD 2026  ·  Alicante  ·  16 de junio de 2026",
      Cm(3.5), Cm(8.5), Cm(26.4), Cm(0.85),
      size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Número de diapositiva (no en portada)


# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 2 — MOTIVACIÓN
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "Mejora Genética del Software — ¿Por qué?")
footer(sl)

items = [
    {"text": "El software contiene comportamientos no óptimos en propiedades no funcionales:\n    tiempo de ejecución, memoria, consumo energético", "bold": False},
    {"text": "Idea: tratar el código como material genético que puede mutarse y evolucionar\n    automáticamente (Petke et al., 2018)", "bold": False},
    {"text": "Proceso: generar variante → compilar → pasar tests funcionales → medir fitness", "bold": False},
    {"text": "Aplicaciones: reparación automática de bugs (APR), optimización de rendimiento", "bold": False},
    {"text": "Objetivo de este trabajo: reducción del tiempo de ejecución\n    manteniendo la corrección funcional del software original", "bold": True, "color": NAVY},
]
bullet_frame(sl, items, Cm(1.0), Cm(2.5), W - Cm(2.0), H - Cm(3.8), size=22)

# Diagrama simplificado
rect(sl, Cm(2), Cm(11.5), Cm(4.5), Cm(1.1), BLUE2)
txbox(sl, "Programa\noriginal", Cm(2), Cm(11.5), Cm(4.5), Cm(1.1),
      size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "→", Cm(6.7), Cm(11.55), Cm(1.2), Cm(0.9), size=22, color=NAVY, align=PP_ALIGN.CENTER)
rect(sl, Cm(8.0), Cm(11.5), Cm(4.5), Cm(1.1), BLUE2)
txbox(sl, "Variante\n(patch)", Cm(8.0), Cm(11.5), Cm(4.5), Cm(1.1),
      size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "→", Cm(12.7), Cm(11.55), Cm(1.2), Cm(0.9), size=22, color=NAVY, align=PP_ALIGN.CENTER)
rect(sl, Cm(14.0), Cm(11.5), Cm(4.5), Cm(1.1), BLUE2)
txbox(sl, "Compilar\n+ Tests", Cm(14.0), Cm(11.5), Cm(4.5), Cm(1.1),
      size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "→", Cm(18.7), Cm(11.55), Cm(1.2), Cm(0.9), size=22, color=NAVY, align=PP_ALIGN.CENTER)
rect(sl, Cm(20.0), Cm(11.5), Cm(4.5), Cm(1.1), GREEN)
txbox(sl, "Medir\nfitness", Cm(20.0), Cm(11.5), Cm(4.5), Cm(1.1),
      size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "→", Cm(24.7), Cm(11.55), Cm(1.2), Cm(0.9), size=22, color=NAVY, align=PP_ALIGN.CENTER)
rect(sl, Cm(26.0), Cm(11.5), Cm(4.5), Cm(1.1), ORANGE)
txbox(sl, "Mejor\nvariante", Cm(26.0), Cm(11.5), Cm(4.5), Cm(1.1),
      size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 3 — ESTADO DEL ARTE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "Metaheurísticas en GI — El debate abierto")
footer(sl)

items = [
    {"text": "Programación Genética (GP): estrategia dominante históricamente\n    → cruces y mutaciones en el espacio de variantes", "bold": False},
    {"text": "Búsqueda Local (LS): hill climbing first-improvement\n    → más efectiva y eficiente que GP (Blot & Petke, 2021, 18 estrategias)", "bold": False},
    {"text": "Búsqueda Aleatoria (RS): sorprendentemente competitiva a veces\n    → evidencia la complejidad del espacio de búsqueda", "bold": False},
    {"text": "MAGPIE: framework unificado para comparar algoritmos de GI\n    (representación estándar de edits, cross-validation, reproducibilidad)", "bold": False},
    {"text": "Pregunta abierta: ¿existe una metaheurística robusta y general\n    para la Mejora Genética?", "bold": True, "color": ORANGE},
]
bullet_frame(sl, items, Cm(1.0), Cm(2.5), W - Cm(2.0), H - Cm(3.8), size=22)


# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 4 — HIPÓTESIS: ¿POR QUÉ UN EDA?
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "Hipótesis — ¿Por qué un Algoritmo de Estimación de Distribuciones?")
footer(sl)

items = [
    {"text": "El espacio de variantes en GI está dominado por variantes inválidas\n    (no compilan o no pasan los tests funcionales)", "bold": False},
    {"text": "GP: saltos globales ciegos → alta probabilidad de caer en zonas inválidas", "bold": False},
    {"text": "LS: avance local seguro → riesgo de quedar atrapada en mínimos locales", "bold": False},
    {"text": "EDA (UMDA): aprende distribución marginal sobre los edits que producen\n    variantes válidas → saltos globales informados",
     "bold": True, "color": NAVY},
    {"text": "Ventaja adicional: sin hiperparámetros de cruce/mutación\n    → configuración simple y resultados más generalizables", "bold": False},
]
bullet_frame(sl, items, Cm(1.0), Cm(2.5), Cm(20), H - Cm(3.8), size=22)

# Diagrama comparativo (columnas)
for xi, (label, desc, col) in enumerate([
    ("GP",   "Saltos\naleatorios\nglobales",     BLUE2),
    ("LS",   "Pasos\nlocales\nseguros",           BLUE2),
    ("EDA",  "Saltos\ninformados\npor distribución", GREEN),
]):
    cx = Cm(22.0) + Cm(xi * 5.5)
    rect(sl, cx, Cm(3.2), Cm(5.0), Cm(1.0), col)
    txbox(sl, label, cx, Cm(3.2), Cm(5.0), Cm(1.0),
          size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, cx, Cm(4.3), Cm(5.0), Cm(4.0),
         GREEN if label == "EDA" else LGRAY,
         line_color=col)
    txbox(sl, desc, cx + Cm(0.2), Cm(4.5), Cm(4.6), Cm(3.8),
          size=17, color=DGRAY if label != "EDA" else NAVY,
          bold=(label == "EDA"), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 5 — REPRESENTACIÓN DEL INDIVIDUO
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "UMDA4GI — Representación del individuo")
footer(sl)

# Edit definition box
rect(sl, Cm(1.0), Cm(2.5), Cm(20), Cm(2.2), LGRAY, line_color=BLUE2)
txbox(sl, "Un edit =  < Operación,  Ubicación,  [Ingrediente] >",
      Cm(1.5), Cm(2.7), Cm(19), Cm(0.85),
      size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
txbox(sl, "Operaciones: Delete · Swap · Copy · Replace   (nivel de sentencia, vía AST)",
      Cm(1.5), Cm(3.6), Cm(19), Cm(0.8),
      size=18, color=DGRAY, align=PP_ALIGN.CENTER)

# Patch definition
txbox(sl, "Un individuo (parche δ) = secuencia de edits:",
      Cm(1.0), Cm(5.1), Cm(20), Cm(0.8), size=22, bold=False, color=DGRAY)

# Example patch
rect(sl, Cm(1.0), Cm(5.9), Cm(20), Cm(2.5), LGRAY, line_color=BLUE2)
txbox(sl, "Ejemplo:    δ = ( e₁, e₂ )",
      Cm(1.5), Cm(6.1), Cm(19), Cm(0.7),
      size=20, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
txbox(sl, "e₁ = < Swap, s₃, s₄ >   →   intercambia sentencias en líneas 3 y 4",
      Cm(2.0), Cm(6.8), Cm(18), Cm(0.65),
      size=18, color=DGRAY)
txbox(sl, "e₂ = < Delete, s₆ >        →   elimina la sentencia en línea 6",
      Cm(2.0), Cm(7.4), Cm(18), Cm(0.65),
      size=18, color=DGRAY)

# Key decision
rect(sl, Cm(1.0), Cm(8.8), Cm(20), Cm(1.8), RGBColor(0xFF, 0xF0, 0xE0), line_color=ORANGE)
txbox(sl, "⭐  Decisión clave: longitud máxima del parche kmax = 2 edits",
      Cm(1.4), Cm(8.95), Cm(19), Cm(0.75),
      size=21, bold=True, color=ORANGE)
txbox(sl, "Parches cortos → más variantes válidas · menor sobreajuste · mejor generalización",
      Cm(1.4), Cm(9.7), Cm(19), Cm(0.7),
      size=18, color=DGRAY)

# Right panel — three operators illustrated
for yi, (op, ex) in enumerate([
    ("Delete s₆",         "Elimina sentencia en línea 6"),
    ("Swap s₃ ↔ s₄",     "Intercambia líneas 3 y 4"),
    ("Replace s₂ ← s₇",  "Sustituye línea 2 por línea 7"),
]):
    ry = Cm(2.5) + yi * Cm(2.8)
    rect(sl, Cm(22), ry, Cm(11.5), Cm(2.4), LGRAY, line_color=BLUE2)
    txbox(sl, op, Cm(22.3), ry + Cm(0.15), Cm(11), Cm(0.9),
          size=20, bold=True, color=NAVY)
    txbox(sl, ex, Cm(22.3), ry + Cm(1.0), Cm(11), Cm(0.8),
          size=17, color=DGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 6 — ALGORITMO UMDA4GI
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "UMDA4GI — Algoritmo")
footer(sl)

# Pseudocode box
rect(sl, Cm(0.8), Cm(2.4), Cm(17.5), Cm(8.4), RGBColor(0xF8,0xF8,0xF8),
     line_color=RGBColor(0xCC,0xCC,0xCC))
pseudo_lines = [
    ("Inicializar  P₀  aleatoriamente  (1 edit / individuo)",     False, DGRAY),
    ("Evaluar P₀  →  identificar mejor δ*",                       False, DGRAY),
    ("",                                                          False, DGRAY),
    ("while  no agotado presupuesto  do",                         True,  NAVY),
    ("    Psel  ←  variantes válidas de Pᵢ  (compilan + tests)", False, DGRAY),
    ("    Aprender distribución  Mᵢ  sobre edits en Psel",        False, BLUE2),
    ("    Pnew  ←  { δ* }  (elitismo)",                           False, DGRAY),
    ("    while  |Pnew| < N - 1  do",                             True,  NAVY),
    ("        Muestrear nuevo individuo δnew  ~  Mᵢ",             False, DGRAY),
    ("        Si δnew  ∈  Pnew  →  sustituir por individuo aleatorio",
                                                                  False, RGBColor(0x80,0,0)),
    ("        Pnew  ←  Pnew  ∪  { δnew }",                       False, DGRAY),
    ("    end while",                                             True,  NAVY),
    ("    Evaluar Pnew;   Pᵢ₊₁ ← Pnew",                          False, DGRAY),
    ("end while  →  devolver mejor individuo",                    True,  NAVY),
]
tb = slide.shapes.add_textbox(Cm(1.2), Cm(2.6), Cm(17), Cm(8.2)) \
    if False else sl.shapes.add_textbox(Cm(1.2), Cm(2.6), Cm(17), Cm(8.2))
tf  = tb.text_frame
tf.word_wrap = True
first = True
for (text, bold, color) in pseudo_lines:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(2)
    run = p.add_run()
    run.text = text
    run.font.name  = "Courier New"
    run.font.size  = Pt(16)
    run.font.bold  = bold
    run.font.color.rgb = color

# Right panel — key insights
for yi, (icon, title, desc) in enumerate([
    ("🔵", "Psel dinámico",
     "Se aprende de TODOS los\nindividuos válidos, no un %\nfijo de la población"),
    ("🟠", "Diversidad garantizada",
     "Individuo duplicado →\nreemplazado por uno\naleatorio (nuevos edits)"),
    ("🟢", "Sin hiperparámetros\nde cruce/mutación",
     "Solo N y kmax=2\nnecesarios"),
]):
    ry = Cm(2.5) + yi * Cm(3.0)
    rect(sl, Cm(19.0), ry, Cm(14.8), Cm(2.7), LGRAY, line_color=BLUE2)
    txbox(sl, title, Cm(19.4), ry + Cm(0.15), Cm(14), Cm(0.85),
          size=19, bold=True, color=NAVY)
    txbox(sl, desc, Cm(19.4), ry + Cm(0.9), Cm(14), Cm(1.6),
          size=17, color=DGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 7 — METODOLOGÍA EXPERIMENTAL
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "Diseño experimental")
footer(sl)

# Left: 4 algoritmos
rect(sl, Cm(0.8), Cm(2.4), Cm(16.5), Cm(8.4), LGRAY, line_color=BLUE2)
txbox(sl, "4 Algoritmos comparados", Cm(1.2), Cm(2.5), Cm(15.5), Cm(0.85),
      size=20, bold=True, color=NAVY)
for yi, (alg, desc) in enumerate([
    ("First Improvement (FI)",      "Búsqueda local hill-climbing"),
    ("Genetic Programming (GP)",    "Algoritmo genético clásico (MAGPIE default)"),
    ("Random Search (RS)",          "Baseline aleatoria"),
    ("UMDA4GI  ★",                 "Nuestra propuesta  [pop_size = 50]"),
]):
    ry = Cm(3.5) + yi * Cm(1.7)
    col = NAVY if "★" in alg else DGRAY
    txbox(sl, f"▸ {alg}", Cm(1.4), ry, Cm(15), Cm(0.75), size=19, bold="★" in alg, color=col)
    txbox(sl, f"   {desc}", Cm(1.4), ry + Cm(0.7), Cm(15), Cm(0.75), size=16, color=DGRAY)

# Middle: Presupuestos
rect(sl, Cm(18.0), Cm(2.4), Cm(7.5), Cm(4.0), LGRAY, line_color=BLUE2)
txbox(sl, "5 Presupuestos (max evals)", Cm(18.3), Cm(2.5), Cm(7), Cm(0.8),
      size=18, bold=True, color=NAVY)
txbox(sl, "100 · 250 · 500 · 750 · 1000",
      Cm(18.3), Cm(3.4), Cm(7), Cm(2.5),
      size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# Right: Validación
rect(sl, Cm(26.0), Cm(2.4), Cm(7.5), Cm(4.0), LGRAY, line_color=BLUE2)
txbox(sl, "Validación cruzada", Cm(26.3), Cm(2.5), Cm(7), Cm(0.8),
      size=18, bold=True, color=NAVY)
txbox(sl, "5-fold\ncross-validation",
      Cm(26.3), Cm(3.4), Cm(7), Cm(2.5),
      size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# Bottom row
rect(sl, Cm(18.0), Cm(7.0), Cm(15.5), Cm(3.6), LGRAY, line_color=BLUE2)
txbox(sl, "Infraestructura HPC", Cm(18.3), Cm(7.1), Cm(15), Cm(0.8),
      size=18, bold=True, color=NAVY)
txbox(sl, "Clúster galgo (SLURM) · Nodo dedicado por experimento\nMétrica: wall-clock time (perf_counter Python)",
      Cm(18.3), Cm(7.9), Cm(15), Cm(2.3), size=17, color=DGRAY)

txbox(sl, "⚠  Total: 4 alg × 4 sw × 5 budgets × 5 folds = 400 experimentos",
      Cm(0.8), Cm(11.1), W - Cm(1.6), Cm(0.9),
      size=20, bold=True, color=ORANGE)


# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 8 — BANCO DE SOFTWARE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "Banco de Software")
footer(sl)

txbox(sl, "4 aplicaciones reales de código abierto · 3 lenguajes · dominios distintos",
      Cm(1.0), Cm(2.4), W - Cm(2.0), Cm(0.75),
      size=20, italic=True, color=DGRAY)

# Software cards
card_data = [
    ("MiniSAT-hack",  "C++",  "SAT Solving",          "Solver.cc",
     "20 archivos .cnf  (10 SAT + 10 UNSAT)\n50–200 variables"),
    ("Sat4J",         "Java", "SAT Solving",          "(múltiples clases)",
     "Mismo corpus que MiniSAT-hack"),
    ("Weka",         "Java", "Machine Learning\n(Random Forest)",  "RandomForest.java",
     "20 datasets .arff  (adjuntos al SW)"),
    ("OptimPNG",      "C",    "Optimización PNG",     "optipng.c, optmi.c",
     "29 imágenes PNG  (color, repositorio ImageMagick)"),
]
colors = [NAVY, BLUE2, ORANGE, GREEN]
for ci, (name, lang, domain, target, data) in enumerate(card_data):
    cx = Cm(0.8) + ci * Cm(8.4)
    # card background
    rect(sl, cx, Cm(3.3), Cm(8.0), Cm(7.8), LGRAY, line_color=colors[ci])
    # header
    rect(sl, cx, Cm(3.3), Cm(8.0), Cm(1.1), colors[ci])
    txbox(sl, name, cx + Cm(0.2), Cm(3.35), Cm(7.6), Cm(1.0),
          size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # content
    for ri, (label, val) in enumerate([
        ("Lenguaje:",   lang),
        ("Dominio:",    domain),
        ("Objetivo:",   target),
        ("Datos:",      data),
    ]):
        ry = Cm(4.6) + ri * Cm(1.55)
        txbox(sl, label, cx + Cm(0.3), ry, Cm(2.8), Cm(0.6),
              size=15, bold=True, color=colors[ci])
        txbox(sl, val,   cx + Cm(3.1), ry, Cm(4.6), Cm(0.9),
              size=15, color=DGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 9 — RESULTADOS I: FITNESS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "Resultados I — Reducción de tiempo de ejecución (validación)")
footer(sl)

# Figura principal
png_bytes = svg_to_png_bytes(os.path.join(PLOTS, "ReductionTime_2x2.svg"), scale=2)
insert_image_bytes(sl, png_bytes, Cm(0.5), Cm(2.3), width=Cm(22.5))

# Ranking table
txbox(sl, "Ranking global (promedio sobre todos los presupuestos):",
      Cm(23.5), Cm(2.3), Cm(10.5), Cm(0.7),
      size=17, bold=True, color=NAVY)

headers = ["Algoritmo", "minisat", "sat4j", "weka", "pngOptim"]
rows = [
    ["First Imp.", "1", "4", "2", "4"],
    ["Gen. Prog.", "3", "2", "1", "2"],
    ["Rand. S.",   "4", "3", "4", "3"],
    ["UMDA4GI",   "2", "1", "3", "1"],
]
col_widths = [Cm(3.5), Cm(1.7), Cm(1.7), Cm(1.7), Cm(1.8)]
table_slide(sl, headers, rows,
            left=Cm(23.5), top=Cm(3.1),
            col_widths=col_widths, row_height=Cm(0.85),
            font_size=15)

# Key message
rect(sl, Cm(23.5), Cm(7.2), Cm(10.5), Cm(2.5), RGBColor(0xE6,0xF0,0xFF), line_color=NAVY)
txbox(sl, "UMDA4GI: único algoritmo con\n2 primeros puestos\n(sat4j y pngOptim)",
      Cm(23.7), Cm(7.35), Cm(10.1), Cm(2.2),
      size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

txbox(sl, "↑ Menor es mejor (reducción %)  ·  Eje Y invertido en figura",
      Cm(0.5), Cm(11.2), Cm(22.5), Cm(0.6),
      size=13, italic=True, color=DGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 10 — RESULTADOS II: SOBREAJUSTE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "Resultados II — Sobreajuste (overfitting)")
footer(sl)

# 2x2 de overfitting plots
ov_files = [
    ("minisat_hack-overfitting.svg", "MiniSAT-hack"),
    ("sat4j-overfitting.svg",        "Sat4J"),
    ("pngOptim-overfitting.svg",     "OptimPNG"),
    ("weka-overfitting.svg",         "Weka"),
]
positions = [
    (Cm(0.4),  Cm(2.3)),
    (Cm(0.4),  Cm(7.6)),
    (Cm(10.5), Cm(2.3)),
    (Cm(10.5), Cm(7.6)),
]
for (fname, title), (lx, ly) in zip(ov_files, positions):
    png_b = svg_to_png_bytes(os.path.join(PLOTS, fname), scale=1.5)
    insert_image_bytes(sl, png_b, lx, ly, width=Cm(10.0))

# Ranking
txbox(sl, "Ranking — menor sobreajuste\n(sin RandomSearch):",
      Cm(21.3), Cm(2.3), Cm(12.5), Cm(0.9),
      size=17, bold=True, color=NAVY)

headers2 = ["Algoritmo", "minisat", "sat4j", "weka", "pngOptim"]
rows2 = [
    ["First Imp.", "2", "3", "3", "3"],
    ["Gen. Prog.", "3", "2", "1", "2"],
    ["UMDA4GI",   "1", "1", "2", "1"],
]
col_widths2 = [Cm(3.5), Cm(1.7), Cm(1.7), Cm(1.7), Cm(1.8)]
table_slide(sl, headers2, rows2,
            left=Cm(21.3), top=Cm(3.3),
            col_widths=col_widths2, row_height=Cm(0.85),
            font_size=15)

rect(sl, Cm(21.3), Cm(6.2), Cm(12.5), Cm(3.0),
     RGBColor(0xE6,0xF0,0xFF), line_color=NAVY)
txbox(sl, "UMDA4GI: menor sobreajuste\nen 3 de 4 softwares\n(2º en Weka)",
      Cm(21.5), Cm(6.35), Cm(12.1), Cm(2.7),
      size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

txbox(sl, "Líneas sólidas = validación · Líneas discontinuas = entrenamiento",
      Cm(0.4), Cm(14.9), Cm(20.5), Cm(0.6),
      size=13, italic=True, color=DGRAY)

rect(sl, Cm(21.3), Cm(9.6), Cm(12.5), Cm(1.5), LGRAY, line_color=BLUE2)
txbox(sl, "Causa probable: kmax=2 edits\n→ parches cortos generalizan mejor",
      Cm(21.5), Cm(9.7), Cm(12.1), Cm(1.3),
      size=16, color=DGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 11 — RESULTADOS III: TIEMPO + LONGITUD
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "Resultados III — Tiempo de búsqueda y longitud de parches")
footer(sl)

# Tiempo
txbox(sl, "Tiempo de búsqueda (minutos):",
      Cm(0.5), Cm(2.3), Cm(16.5), Cm(0.65),
      size=18, bold=True, color=NAVY)
tiempo_bytes = pdf_to_png_bytes(os.path.join(PLOTS, "tiempo_2x2.pdf"), dpi=180)
insert_image_bytes(sl, tiempo_bytes, Cm(0.5), Cm(3.0), width=Cm(17.0))

# Longitud parches
txbox(sl, "Longitud media del parche (nº de edits):",
      Cm(18.5), Cm(2.3), Cm(15.5), Cm(0.65),
      size=18, bold=True, color=NAVY)
pl_bytes = svg_to_png_bytes(os.path.join(PLOTS, "patch_length_mean_barh.svg"), scale=1.8)
insert_image_bytes(sl, pl_bytes, Cm(18.5), Cm(3.0), width=Cm(15.5))

# Key messages
txbox(sl, "▸ UMDA4GI y GP: los algoritmos más rápidos",
      Cm(0.5), Cm(11.2), Cm(16.5), Cm(0.65), size=17, color=DGRAY)
txbox(sl, "▸ Parches cortos (kmax=2) → el menor número de edits entre los algoritmos poblacionales",
      Cm(0.5), Cm(11.85), Cm(33), Cm(0.65), size=17, color=DGRAY)
txbox(sl, "▸ Parches más largos (FI) → más sobreajuste (consistente con Petke et al., 2023)",
      Cm(0.5), Cm(12.5), Cm(33), Cm(0.65), size=17, color=DGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 12 — CONCLUSIONES
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "Conclusiones")
footer(sl)

conclusions = [
    ("1", "FITNESS COMPETITIVO",
     "UMDA4GI obtiene los mejores resultados en 2 de 4 softwares\n"
     "(sat4j, pngOptim) y posiciones competitivas en los otros dos.\n"
     "Es el único algoritmo con dos primeros puestos en el ranking global.",
     GREEN),
    ("2", "MENOR SOBREAJUSTE",
     "UMDA4GI generaliza mejor que FI y GP en 3 de 4 softwares.\n"
     "Los parches de 2 edits reducen el overfitting al no\n"
     "sobreadaptar la búsqueda a los datos de entrenamiento.",
     BLUE2),
    ("3", "VELOCIDAD INTERMEDIA",
     "Tiempo de búsqueda similar a GP, sin sus parámetros\n"
     "de cruce/mutación. Fácil de configurar y reproducir.",
     ORANGE),
]

for ci, (num, title, desc, col) in enumerate(conclusions):
    cy = Cm(2.5) + ci * Cm(3.4)
    rect(sl, Cm(0.8), cy, Cm(33), Cm(3.2), LGRAY, line_color=col)
    rect(sl, Cm(0.8), cy, Cm(1.8), Cm(3.2), col)
    txbox(sl, num, Cm(0.8), cy + Cm(0.9), Cm(1.8), Cm(1.4),
          size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, title, Cm(2.9), cy + Cm(0.1), Cm(30), Cm(0.85),
          size=20, bold=True, color=col)
    txbox(sl, desc, Cm(2.9), cy + Cm(0.95), Cm(30), Cm(2.1),
          size=18, color=DGRAY)

# Limitación
rect(sl, Cm(0.8), Cm(13.0), Cm(33), Cm(1.5),
     RGBColor(0xFF,0xF0,0xE0), line_color=ORANGE)
txbox(sl, "⚠  Limitación: la propuesta depende de la calidad de la suite de tests\n"
          "   (sin tests suficientes, puede generar variantes funcionalmente incorrectas)",
      Cm(1.2), Cm(13.1), Cm(32.5), Cm(1.3), size=17, color=DGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 13 — TRABAJO FUTURO
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
header_bar(sl, "Trabajo futuro")
footer(sl)

future = [
    ("Filtrado de Psel",
     "Aprender la distribución solo sobre variantes que mejoran el fitness\n"
     "(no sobre todas las válidas) → ¿distribución de mejora más pura?"),
    ("LLMs como operadores (en progreso)",
     "Usar modelos de lenguaje para generar edits semánticamente informados.\n"
     "UMDA4GI + LLM (UMDA4GI): LLM aprende de las mejores variantes → genera\n"
     "los edits de la siguiente generación guiándose por el código y el historial."),
    ("LLM puro como baseline",
     "¿Puede un LLM solo (sin EDA) mejorar el software?\n"
     "Experimentos en ejecución actualmente para aislar la contribución del LLM."),
    ("Análisis semántico de la distribución",
     "¿Qué aprende UMDA sobre el código? ¿Hay edits recurrentes con significado?\n"
     "Interpretabilidad del modelo probabilístico aprendido."),
]

for fi, (title, desc) in enumerate(future):
    fy = Cm(2.5) + fi * Cm(2.8)
    is_llm = "LLM" in title
    col = ORANGE if is_llm else NAVY
    rect(sl, Cm(0.8), fy, Cm(33), Cm(2.6), LGRAY, line_color=col)
    txbox(sl, ("🔬 " if is_llm else "→ ") + title,
          Cm(1.2), fy + Cm(0.1), Cm(32), Cm(0.8),
          size=19, bold=True, color=col)
    txbox(sl, desc, Cm(1.5), fy + Cm(0.9), Cm(32), Cm(1.5),
          size=17, color=DGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 14 — CIERRE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl, NAVY)
rect(sl, 0, Cm(5.5), W, Cm(0.25), GOLD)

txbox(sl, "¡Gracias!  /  ¿Preguntas?",
      Cm(1.5), Cm(1.5), W - Cm(3.0), Cm(2.0),
      size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Contact info boxes
for xi, (icon, label, val) in enumerate([
    ("✉", "Email",    "pablo.bermejo@uclm.es"),
    ("🐙", "Código",  "github.com/paberlo/magpie"),
    ("📦", "Datos",   "zenodo.org/records/19469242"),
]):
    cx = Cm(1.5) + xi * Cm(11.5)
    rect(sl, cx, Cm(6.2), Cm(11.0), Cm(2.5), BLUE2)
    txbox(sl, label, cx + Cm(0.3), Cm(6.3), Cm(10.4), Cm(0.75),
          size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txbox(sl, val,   cx + Cm(0.3), Cm(7.1), Cm(10.4), Cm(0.85),
          size=16, color=WHITE, align=PP_ALIGN.CENTER)

txbox(sl, "UMDA aplicado a la Mejora Automática del Software  ·  JISBD 2026  ·  Alicante",
      Cm(1.5), Cm(9.2), W - Cm(3.0), Cm(0.7),
      size=16, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# ── Guardar ───────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"✅  Presentación guardada en: {OUT}")
print(f"    Diapositivas: {len(prs.slides)}")
